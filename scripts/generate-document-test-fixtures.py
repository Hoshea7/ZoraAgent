"""Generate representative Office/PDF fixtures for document-reader integration tests.

Run with the bundled Codex Python runtime documented by the desktop app. The
generated binaries are committed so normal test runs do not require Python or
the generation libraries.
"""

from datetime import date
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Font, PatternFill
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from pypdf import PdfReader
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "tests" / "fixtures" / "documents"


def generate_docx() -> None:
    document = Document()
    section = document.sections[0]
    section.page_width = Inches(8.27)
    section.page_height = Inches(11.69)
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.8)

    document.core_properties.title = "Northstar Product Review"
    document.core_properties.author = "Zora QA"
    title = document.add_heading("Northstar Product Review", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_paragraph("Reporting period: Q2 2026")

    document.add_heading("Executive summary", level=1)
    document.add_paragraph(
        "The Northstar pilot reached 1,248 active teams. Weekly document reads "
        "increased by 18 percent while median processing time remained below two seconds."
    )
    document.add_paragraph("DOCX_REALISTIC_MARKER")

    document.add_heading("Release priorities", level=2)
    for item in [
        "Preserve heading and list structure during extraction",
        "Return tables as readable Markdown",
        "Keep internal attachment paths out of model-visible content",
    ]:
        document.add_paragraph(item, style="List Bullet")

    table = document.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    headers = ["Workstream", "Owner", "Status"]
    for cell, value in zip(table.rows[0].cells, headers):
        cell.text = value
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for run in cell.paragraphs[0].runs:
            run.bold = True
    for row in [
        ("Document parser", "Mei", "Complete"),
        ("Runtime integration", "Arun", "In review"),
        ("E2E coverage", "Sofia", "In progress"),
    ]:
        cells = table.add_row().cells
        for cell, value in zip(cells, row):
            cell.text = value

    document.add_section(WD_SECTION.NEW_PAGE)
    document.add_heading("Appendix", level=1)
    document.add_paragraph("Validation checklist", style="List Number")
    document.add_paragraph("Open the file in Microsoft Word or LibreOffice.", style="List Number")
    document.add_paragraph("Confirm that the table and page break are retained.", style="List Number")
    footer = document.sections[-1].footer.paragraphs[0]
    footer.text = "Confidential test fixture"
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in footer.runs:
        run.font.size = Pt(9)

    target = OUTPUT / "northstar-review.docx"
    document.save(target)
    reopened = Document(target)
    assert len(reopened.tables) == 1
    assert len(reopened.paragraphs) >= 10


def generate_xlsx() -> None:
    workbook = Workbook()
    summary = workbook.active
    summary.title = "Summary"
    summary.merge_cells("A1:D1")
    summary["A1"] = "Northstar Usage Dashboard"
    summary["A1"].font = Font(name="Arial", size=16, bold=True, color="FFFFFF")
    summary["A1"].fill = PatternFill("solid", fgColor="1F4E78")
    summary["A1"].alignment = Alignment(horizontal="center")
    summary.append(["Month", "Active teams", "Document reads", "Success rate"])
    for cell in summary[2]:
        cell.font = Font(name="Arial", bold=True)
        cell.fill = PatternFill("solid", fgColor="D9EAF7")

    rows = [
        (date(2026, 4, 1), 980, 18240, 0.972),
        (date(2026, 5, 1), 1105, 21480, 0.981),
        (date(2026, 6, 1), 1248, 25350, 0.986),
    ]
    for row in rows:
        summary.append(row)
    summary["A6"] = "Quarter total"
    summary["B6"] = "=MAX(B3:B5)"
    summary["C6"] = "=SUM(C3:C5)"
    summary["D6"] = "=AVERAGE(D3:D5)"
    summary["A8"] = "XLSX_REALISTIC_MARKER"
    for cell in summary[6]:
        cell.font = Font(name="Arial", bold=True)
    for row in range(3, 6):
        summary[f"A{row}"].number_format = "yyyy-mm-dd"
        summary[f"D{row}"].number_format = "0.0%"
    summary["D6"].number_format = "0.0%"
    summary.freeze_panes = "A3"
    summary.auto_filter.ref = "A2:D5"
    for column, width in {"A": 22, "B": 16, "C": 18, "D": 16}.items():
        summary.column_dimensions[column].width = width

    data = workbook.create_sheet("Regional Data")
    data.append(["Region", "Teams", "Reads"])
    data.append(["East", 460, 9200])
    data.append(["Central", 388, 8100])
    data.append(["West", 400, 8050])
    notes = workbook.create_sheet("Internal Notes")
    notes.sheet_state = "hidden"
    notes["A1"] = "Generated fixture: do not use as business data"

    workbook.calculation.fullCalcOnLoad = True
    workbook.calculation.forceFullCalc = True
    workbook.calculation.calcMode = "auto"
    target = OUTPUT / "northstar-dashboard.xlsx"
    workbook.save(target)
    reopened = load_workbook(target, data_only=False)
    assert reopened["Summary"]["C6"].value == "=SUM(C3:C5)"
    assert reopened["Internal Notes"].sheet_state == "hidden"


def generate_pptx() -> None:
    presentation = Presentation()
    presentation.slide_width = PptxInches(13.333)
    presentation.slide_height = PptxInches(7.5)
    presentation.core_properties.title = "Northstar Launch Readout"
    presentation.core_properties.author = "Zora QA"

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Northstar Launch Readout"
    title_slide.placeholders[1].text = "Q2 2026 · Document intelligence"
    title_slide.notes_slide.notes_text_frame.text = (
        "Open with the adoption result and clarify that all figures are test data. "
        "PPTX_SPEAKER_NOTE_MARKER"
    )

    summary_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    summary_slide.shapes.title.text = "Adoption and reliability"
    chart_data = ChartData()
    chart_data.categories = ["April", "May", "June"]
    chart_data.add_series("Active teams", (980, 1105, 1248))
    summary_slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        PptxInches(0.7),
        PptxInches(1.4),
        PptxInches(6.2),
        PptxInches(4.8),
        chart_data,
    )
    text_box = summary_slide.shapes.add_textbox(
        PptxInches(7.15), PptxInches(1.65), PptxInches(5.55), PptxInches(3.6)
    )
    frame = text_box.text_frame
    frame.text = "PPTX_REALISTIC_MARKER"
    for item in [
        "1,248 active teams in June",
        "98.6 percent successful reads",
        "Median processing time below two seconds",
    ]:
        paragraph = frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = PptxPt(18)
    frame.paragraphs[0].font.bold = True
    frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    summary_slide.notes_slide.notes_text_frame.text = "Call out the June value and the reliability trend."

    table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_slide.shapes.title.text = "Release readiness"
    table_data = [
        ["Capability", "Owner", "Status"],
        ["PDF and Office parsing", "Mei", "Ready"],
        ["Claude runtime", "Arun", "Ready"],
        ["Pi runtime", "Sofia", "Ready"],
    ]
    table_slide.shapes.add_table(
        len(table_data),
        len(table_data[0]),
        PptxInches(0.8),
        PptxInches(1.5),
        PptxInches(11.7),
        PptxInches(3.4),
    )
    table = table_slide.shapes[-1].table
    for row_index, row in enumerate(table_data):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptxPt(16)
                paragraph.font.bold = row_index == 0

    target = OUTPUT / "northstar-launch.pptx"
    presentation.save(target)
    reopened = Presentation(target)
    assert len(reopened.slides) == 3
    assert "PPTX_SPEAKER_NOTE_MARKER" in reopened.slides[0].notes_slide.notes_text_frame.text


def generate_pdf() -> None:
    target = OUTPUT / "northstar-operations.pdf"
    styles = getSampleStyleSheet()
    document = SimpleDocTemplate(
        str(target),
        pagesize=A4,
        rightMargin=20 * mm,
        leftMargin=20 * mm,
        topMargin=18 * mm,
        bottomMargin=18 * mm,
        title="Northstar Operations Report",
        author="Zora QA",
    )
    story = [
        Paragraph("Northstar Operations Report", styles["Title"]),
        Spacer(1, 8 * mm),
        Paragraph("PDF_REALISTIC_MARKER", styles["Heading2"]),
        Paragraph(
            "This representative fixture contains multiple pages, headings, paragraphs, "
            "a table, repeated headers, and document metadata.",
            styles["BodyText"],
        ),
        Spacer(1, 6 * mm),
        Table(
            [
                ["Metric", "April", "May", "June"],
                ["Active teams", "980", "1,105", "1,248"],
                ["Document reads", "18,240", "21,480", "25,350"],
                ["Success rate", "97.2%", "98.1%", "98.6%"],
            ],
            colWidths=[48 * mm, 32 * mm, 32 * mm, 32 * mm],
            repeatRows=1,
            style=TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1F4E78")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
                    ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ]
            ),
        ),
        PageBreak(),
        Paragraph("Operational notes", styles["Heading1"]),
        Paragraph(
            "The second page verifies page boundaries and stable pagination. "
            "All names and metrics in this file are synthetic test data.",
            styles["BodyText"],
        ),
        Spacer(1, 5 * mm),
        Paragraph("Validation owner: Document Platform QA", styles["BodyText"]),
    ]
    document.build(story)
    reopened = PdfReader(target)
    assert len(reopened.pages) == 2
    assert reopened.metadata.title == "Northstar Operations Report"


def main() -> None:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    generate_docx()
    generate_xlsx()
    generate_pptx()
    generate_pdf()


if __name__ == "__main__":
    main()
